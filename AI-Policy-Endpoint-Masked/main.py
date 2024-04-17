# ------------------------------------------------- #
#   IMPORTING MODULES
# ------------------------------------------------- #
import asyncio
import shutil
from enum import Enum
from typing import Optional, List, Annotated, Dict, Any

import PyPDF2
import pandas as pd
from fastapi import FastAPI, HTTPException, Query, Form
from pydantic import BaseModel, Field
from fastapi import FastAPI, File, UploadFile, Body
from fastapi.responses import JSONResponse

# import warnings
# warnings.filterwarnings("ignore")

from utils import cleanup_text, read_files, sort_dict_by_key, extract_keywords, scrape_imf_search_results, \
    scrape_worldbank_search_results, detect_country

import openai
import tiktoken
from pprint import pprint
from getpass import getpass

from pydantic.dataclasses import dataclass
from rich.markdown import Markdown

import os
import glob
import textwrap
import time
import re
import langchain

# callbacks openai
from langchain_community.callbacks import get_openai_callback

# loaders
from langchain_community.document_loaders import DirectoryLoader

# splits
from langchain.text_splitter import RecursiveCharacterTextSplitter

# prompts
# from langchain import PromptTemplate, ConversationChain, LLMChain
from langchain.chains import ConversationChain, LLMChain

# vector stores
from langchain_community.vectorstores import Chroma

# retrievers
from langchain.chains import RetrievalQA, ConversationalRetrievalChain

# json parser
from langchain_core.output_parsers import JsonOutputParser

# from langchain_community.embeddings import OpenAIEmbeddings
from langchain_openai import OpenAIEmbeddings
from langchain.prompts import PromptTemplate

from langchain.chains import RetrievalQA
from langchain_openai import OpenAI

from dotenv import load_dotenv

from langchain.chains import RetrievalQA
from langchain_openai import OpenAI
from langchain_openai import ChatOpenAI

# from langchain.chat_models import ChatOpenAI
import tiktoken
import fitz
from docx import Document
from pptx import Presentation
import openpyxl
import io
from masked_ai.masker import Masker
from rouge_score import rouge_scorer
from transformers import BertTokenizer, BertModel
from bert_score import BERTScorer

# summarization modules
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.probability import FreqDist
from nltk.tokenize.treebank import TreebankWordDetokenizer
import nltk
from langchain.chains.openai_functions import create_structured_output_chain
from langchain.prompts import ChatPromptTemplate

# Download the 'punkt' resource
nltk.download('punkt')
# Download the stopwords corpus
nltk.download('stopwords')

scorer = rouge_scorer.RougeScorer(['rougeL'], use_stemmer=True)

# ------------------------------------------------- #
#  LOAD ENV VARIABEL & VECTOR
# ------------------------------------------------- #
load_dotenv()
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
EMBEDDING_PERSIST_DIR = os.getenv('EMBEDDING_PERSIST_DIR')
MODEL_NAME = os.getenv('MODEL_NAME')
MAX_TOKENS_COMPLETION_FOR_GENERATE_OUTLINE = int(os.getenv('MAX_TOKENS_COMPLETION_FOR_GENERATE_OUTLINE'))
MAX_TOKENS_COMPLETION_FOR_GENERATE_SUBSECTION = int(os.getenv('MAX_TOKENS_COMPLETION_FOR_GENERATE_SUBSECTION'))
MAX_TOKENS_COMPLETION_FOR_GENERATE_CONCEPT_NOTE = int(os.getenv('MAX_TOKENS_COMPLETION_FOR_GENERATE_CONCEPT_NOTE'))
MAX_TOKENS_COMPLETION_FOR_REGENERATE_SUBSECTION = int(os.getenv('MAX_TOKENS_COMPLETION_FOR_REGENERATE_SUBSECTION'))
MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_OUTLINE = int(
    os.getenv('MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_OUTLINE'))
MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_SUBSECTIONS = int(
    os.getenv('MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_SUBSECTIONS'))

vectordb = Chroma(persist_directory=EMBEDDING_PERSIST_DIR,
                  embedding_function=OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY))
# print(vectordb._collection.count())
tokenizer = tiktoken.encoding_for_model(MODEL_NAME)

# ------------------------------------------------- #
#  LOAD MAIN FUNCTION
# ------------------------------------------------- #

# sample query = "Please generate policy outline for Data Protection and Privacy Policy Implementation in Australia's Industry"
# uvicorn main:app --reload


import logging

logger = logging.getLogger(__name__)
app = FastAPI(debug=True)


def count_tokens(data):
    token_counts = (len(tokenizer.encode(data)))
    return token_counts


def nltk_summarize(text, max_words):
    # Tokenize the text into words and sentences
    words = word_tokenize(text)
    sentences = sent_tokenize(text)

    # Remove stopwords
    # Stopwords are common words that often do not carry significant meaning (e.g., "and", "the", "is").
    # The code creates a set of English stopwords and then removes them from the list of words.
    # Additionally, it converts each word to lowercase and filters out non-alphanumeric characters.
    stop_words = set(stopwords.words('english'))
    words = [word.lower() for word in words if word.isalnum() and word.lower() not in stop_words]

    # counts the occurrences of each unique word in the processed text.
    freq_dist = FreqDist(words)

    # Each sentence is scored based on the sum of the frequencies of its constituent words.
    # The higher the sum, the more important the sentence is considered.
    sentence_scores = {sentence: sum([freq_dist[word] for word in word_tokenize(sentence)]) for sentence in sentences}

    # The sentences are sorted in descending order based on their scores. The code then iterates through the sentences,
    # selecting them until the total word count in the selected sentences reaches or exceeds the specified max_words.
    selected_sentences = []
    total_words = 0
    for sentence in sorted(sentence_scores, key=sentence_scores.get, reverse=True):
        if total_words + len(word_tokenize(sentence)) <= max_words:
            selected_sentences.append(sentence)
            total_words += len(word_tokenize(sentence))
        else:
            break

    # Detokenize to form the final summary
    summary = TreebankWordDetokenizer().detokenize(selected_sentences)

    return summary


def get_title(user_prompt, question, country_name):
    # Get title
    context = 'Get a title for the document'
    title_prompt = f"""
            ***** {question}:
            {user_prompt}
            -------------------
            *****
            ***** {country_name}. 
            *****
            *****
            """

    title_chain = RetrievalQA.from_chain_type(
        llm=ChatOpenAI(
            model_name=MODEL_NAME,
            openai_api_key=OPENAI_API_KEY,
            max_tokens=50,
            temperature=0,
        ),
        retriever=vectordb.as_retriever(search_kwargs={'k': 7}),
        return_source_documents=True
    )

    title = title_chain.invoke(title_prompt)['result']
    return title


class PromptImplementation(str, Enum):
    RECOMMENDED_DOCUMENTS = "Recommended Documents"
    TABLE_OF_CONTENT = "Table of Content"
    POLICY_OUTLINE = "Policy Outline"
    SUBSECTIONS_DETAILING = "Subsections Detailing"
    CONCEPT_NOTES = "Concept Notes"
    REGENERATE_SUBSECTION = "Regenerate Subsection"
    REGENERATE_OUTLINES_NOTES = "Regenerate Outlines/Notes"
    REGENERATE_TABLE_OF_CONTENT = "Regenerate Table of Content"


def validate_prompt_input(input_string: str, prompt_implementation: PromptImplementation) -> str:
    required_placeholders = None
    if prompt_implementation == PromptImplementation.POLICY_OUTLINE:
        required_placeholders = ['{user_prompt}', '{country_name}', '{context}', '{table_of_content}']
    elif prompt_implementation == PromptImplementation.SUBSECTIONS_DETAILING:
        required_placeholders = ['{content}', '{subsection_title}', '{section_title}', '{context}', '{country_name}']
    elif prompt_implementation == PromptImplementation.CONCEPT_NOTES:
        required_placeholders = ['{user_prompt}', '{context}', '{country_name}']
    elif prompt_implementation == PromptImplementation.REGENERATE_SUBSECTION:
        required_placeholders = ['{subsection_title}', '{subsection_explanation}', '{user_prompt_for_regenerate}',
                                 '{country_name}']
    elif prompt_implementation == PromptImplementation.REGENERATE_OUTLINES_NOTES:
        required_placeholders = ['{document_format}', '{document_value}', '{country_name}']
    # latest
    elif prompt_implementation == PromptImplementation.RECOMMENDED_DOCUMENTS:
        required_placeholders = ['{user_prompt}', '{doc_type}', '{country_name}']
    elif prompt_implementation == PromptImplementation.TABLE_OF_CONTENT:
        required_placeholders = ['{user_prompt}', '{country_name}', '{context}']
    elif prompt_implementation == PromptImplementation.REGENERATE_TABLE_OF_CONTENT:
        required_placeholders = ['{table_of_content}', '{user_prompt}', '{country_name}']
    else:
        raise HTTPException(status_code=400, detail={"error": "prompt_implementation value required"})

    missing_placeholders = [placeholder for placeholder in required_placeholders if placeholder not in input_string]

    if missing_placeholders:
        raise HTTPException(status_code=400, detail={
            "Missing required placeholders for custom backend prompt": ', '.join(missing_placeholders)})
    if (len(re.findall(r'{\w+}', input_string)) != len(required_placeholders)) or ('{}' in input_string):
        raise HTTPException(status_code=400, detail={
            "error": "Incorrect number of placeholders. Unable to add another variable required."})

    return input_string


def langchain_parsing(policy_outline):
    policy_outline_schema = {
        "title": "Policy Outline",
        "description": "Outline of a policy with sections and subsections.",
        "type": "object",
        "properties": {
            "sections": {
                "title": "Sections",
                "description": "Sections of the policy outline",
                "type": "object",
                "patternProperties": {
                    "^\d+(\.\s.*)?$": {
                        "type": "object",
                        "patternProperties": {
                            "^\d+(\.\s.*)?$": {"type": "string"}
                        }
                    }
                }
            }
        },
        "required": ["sections"]
    }

    llm = ChatOpenAI(model="gpt-3.5-turbo-1106", temperature=0)
    prompt = ChatPromptTemplate.from_messages(
        [
            ("system", "*****"),
            ("human", f""""***** : 
                            {policy_outline}
                            -----------------------------
                            *****
                        """"")
        ]
    )

    chain = create_structured_output_chain(policy_outline_schema, llm, prompt, verbose=True)
    output = chain.run({})
    try:
        parsed_policy_outline = output['policy_outline']
    except:
        parsed_policy_outline = output
    return parsed_policy_outline


def langchain_parsing_for_notes(concept_outline):
    concept_outline_schema = {
        "title": "Concept Notes",
        "description": "A concept notes with sections and subsections.",
        "type": "object",
        "properties": {
            "sections": {
                "title": "Sections",
                "description": "Sections of the concept notes",
                "type": "object",
                "patternProperties": {
                    "^\d+(\.\s.*)?$": {
                        "type": "object",
                        "patternProperties": {
                            "^\d+(\.\s.*)?$": {"type": "string"}
                        }
                    }
                }
            }
        },
        "required": ["sections"]
    }

    llm = ChatOpenAI(model="gpt-3.5-turbo-1106", temperature=0)
    prompt = ChatPromptTemplate.from_messages(
        [
            ("system", "You are extracting information in a structured formats."),
            ("human", f""""***** : 
                            {concept_outline}
                            -----------------------------
                            *****
                        """"")
        ]
    )

    chain = create_structured_output_chain(concept_outline_schema, llm, prompt, verbose=True)
    output = chain.run({})
    try:
        parsed_concept_outline = output['concept_notes']
    except:
        parsed_concept_outline = output
    return parsed_concept_outline


def check_sections(parsed_result):
    try:
        for version in parsed_result['versions']:
            for section in version['sections']:
                if len(section) < 1:
                    return False
    except:
        try:
            for version in parsed_result['versions']:
                for section in version:
                    if len(section) < 1:
                        return False
        except:
            return False
    return True


def toc_parsing(table_of_content):
    table_of_content_schema = {
        "title": "Policy Outline",
        "description": "Outline of a policy with sections and subsections.",
        "type": "object",
        "properties": {
            "versions": {
                "title": "Versions",
                "description": "Different versions of the policy outline",
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "version": {"type": "string"},
                        "sections": {
                            "title": "Sections",
                            "description": "Sections of the policy outline",
                            "type": "object",
                            "patternProperties": {
                                "^\d+(\.\s.*)?$": {
                                    "type": "array",
                                    "items": {"type": "string"}
                                }
                            }
                        }
                    },
                    "required": ["sections"]
                }
            }
        }
    }

    llm = ChatOpenAI(model="gpt-3.5-turbo-1106", temperature=0)
    prompt = ChatPromptTemplate.from_messages(
        [
            ("system", "*****"),
            ("human", f""""***** 
                            {table_of_content}
                            -----------------------------
                            *****
                            *****
                            *****
                            *****
                        """"")
        ]
    )

    chain = create_structured_output_chain(table_of_content_schema, llm, prompt, verbose=True)
    # output = chain.run({})

    retries = 0
    while retries < 3:
        output = chain.run({})
        # Checking
        if check_sections(output):
            parsed_policy_outline = output
            break
        else:
            retries += 1
    else:
        try:
            parsed_policy_outline = chain.run({})['policy_outline']
        except:
            parsed_policy_outline = None

    return parsed_policy_outline


def toc_regenerated_parsing(table_of_content):
    table_of_content_schema_2 = {
        "title": "Table of Content",
        "description": "Table of Content with sections and subsections.",
        "type": "object",
        "properties": {
            "sections": {
                "title": "Sections",
                "description": "Sections of the policy outline",
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "section": {"type": "string"},
                        "subsections": {
                            "type": "array",
                            "items": {"type": "string"}
                        }
                    },
                    "required": ["section", "subsections"]
                }
            }
        }
    }

    llm2 = ChatOpenAI(model="gpt-3.5-turbo-1106", temperature=0)
    prompt2 = ChatPromptTemplate.from_messages(
        [
            ("system", "*****"),
            ("human", f""""*****: 
                            {table_of_content}
                            -----------------------------
                            *****
                            *****
                            *****
                            *****
                        """"")
        ]
    )

    chain2 = create_structured_output_chain(table_of_content_schema_2, llm2, prompt2, verbose=True)
    output = chain2.run({})
    return output


app = FastAPI()


# @TODO : UNHIDE IF NEEDED
# @app.put("/update_backend_prompt_test")
# async def update_backend_prompt(
#     input_string: str,
#     prompt_implementation: PromptImplementation = Query(...)
# ):
#     return validate_prompt_input(input_string, prompt_implementation)


# @TODO : DEBUG SWAGGER FILES, UNHIDE IF NEEDED
# @app.post("/test-debug", status_code=200)
# async def generate_test(owner: str ,
#                         file: List[UploadFile] = File(None)
#                         ):
#     file_sizes = []
#     if file:
#         for fil in file:
#             # contents = await fil.read()
#             file_sizes.append(fil.filename)
#     return {"owner": owner, "file_sizes": file_sizes}


@app.post("/generate_recommended_documents", status_code=200)
async def generate_recommended_documents(
        user_prompt: str,
        country_name: str,
        doc_type: str,
        custom_backend_prompt: Optional[str] = None,
        is_super_admin: Optional[bool] = False,
):
    """
    *****

    *****:
    - *****
    - *****
    - *****
    - *****
         ** ***** **

            *****
            *****
            *****

            *****

            *****
            *****

            ****
            1. *****
            *****
            *****

            2. *****
            *****
            *****

            3. *****
            *****
            *****

            4. *****
            *****
            *****

            5. *****
            *****
            *****

            6. *****
            *****
            *****

            7. *****
            *****
            *****

            8. *****
            *****
            *****
            ****

            *****
            the structure of a concept note are delimited by ####.

            ####
            1.*****:
            *****

            2.*****:
            *****

            3.*****:
            *****

            4.*****:
            *****
            ####

            *****

            *****

            *****

            *****

            *****


    Returns:
    dict: A dictionary containing recommended documents.

    """
    user_prompt = user_prompt
    country_name = country_name
    doc_type = doc_type


    validated_country = country_name
    if country_name.lower() == "other":
        try:
            detected_country_name = detect_country(sentence=user_prompt.lower())
            if detected_country_name is not None:
                validated_country = detected_country_name
            else:
                return JSONResponse(status_code=400, content={"error": "Country name in the user prompt not found"})
        except:
            return JSONResponse(status_code=400, content={"error": "Country name in the user prompt not found"})

    # Admin has different rules
    if is_super_admin == True:
        # default prompt as custom prompt
        default_writing_rule = f"""*****
                                ***** {doc_type} *****,
                                ***** {validated_country} *****.
                                *****
                                *****
                                *****
                                *****
                            
                                1. *****
                                2. *****
                                3. *****
                                4. *****
                                5. *****"""
        recommendation_template = user_prompt + "\n----------\n" + default_writing_rule

        use_default_prompt = False
        recommendation_chain = RetrievalQA.from_chain_type(
            llm=ChatOpenAI(
                model_name=MODEL_NAME,
                openai_api_key=OPENAI_API_KEY,
                max_tokens=500,
                timeout=30,
            ),
            retriever=vectordb.as_retriever(search_kwargs={'k': 7}),
            return_source_documents=True,
        )

        try:
            # Generate the recommended documents
            start_time = time.time()
            total_tokens = 0
            total_prompt_tokens = 0
            total_completion_tokens = 0
            total_cost = 0

            with get_openai_callback() as cb:
                recommended_documents = await recommendation_chain.ainvoke(recommendation_template)

            end_time = time.time()
            time_lapse = f"{end_time - start_time} seconds"

        except Exception as e:
            return JSONResponse(status_code=400, content={"error": str(e)})

        # NON GPT recommendation:
        keywords = extract_keywords(text=user_prompt, country_name=validated_country)
        try:
            scraped_imf = scrape_imf_search_results(keywords)
        except:
            scraped_imf = "Not found"

        try:
            scraped_wordlbank = scrape_worldbank_search_results(keywords)
        except:
            scraped_wordlbank = "Not found"

        result_concated = f"""{recommended_documents["result"]} \n\n<label style="color: #37377d">Recommendations based on online web sources:</label> \nSource: <a href="https://www.imf.org/" target="_blank">https://imf.org/</a>\n{scraped_imf} \n\nSource: <a href="https://www.worldbank.org/" target="_blank">https://www.worldbank.org/</a>\n{scraped_wordlbank}"""

        return {"recommended_documents": result_concated, "prompt": recommendation_template,
                "use_default_prompt": use_default_prompt, "country_name": validated_country, is_super_admin:is_super_admin}

    # For common users
    else:
        use_default_prompt = True
        if custom_backend_prompt:
            prompt_string = validate_prompt_input(custom_backend_prompt, PromptImplementation.RECOMMENDED_DOCUMENTS)
            prompt_string = prompt_string.replace('{user_prompt}', user_prompt)
            prompt_string = prompt_string.replace('{doc_type}', doc_type)
            prompt_string = prompt_string.replace('{country_name}', validated_country)
            recommendation_template = prompt_string
            use_default_prompt = False
        else:
            recommendation_template = f"""
            ***** {user_prompt}.
        
            *****
            *****
        
            ***** : {doc_type}
            
            *****
            ***** ****.
        
            ****
            1. *****
            *****
            *****
        
            2. *****
            *****
            *****
        
            3. *****
            *****
            *****
        
            4. *****
            *****
            *****
        
            5. *****
            *****
            *****
        
            6. *****
            *****
            *****
        
            7. *****
            *****
            *****
        
            8. *****
            *****
            *****
            ****
        
            *****
            the structure of a concept note are delimited by ####.
        
            ####
            1.*****:
            *****
        
            2.*****:
            *****
        
            3.*****:
            *****
        
            4.*****:
            *****
            ####            
            """

        recommendation_chain = RetrievalQA.from_chain_type(
            llm=ChatOpenAI(
                model_name=MODEL_NAME,
                openai_api_key=OPENAI_API_KEY,
                max_tokens=500,
                timeout=30,
            ),
            retriever=vectordb.as_retriever(search_kwargs={'k': 7}),
            return_source_documents=True,
        )

        try:
            # Generate the recommended documents
            start_time = time.time()
            total_tokens = 0
            total_prompt_tokens = 0
            total_completion_tokens = 0
            total_cost = 0

            with get_openai_callback() as cb:
                recommended_documents = await recommendation_chain.ainvoke(recommendation_template)

            end_time = time.time()
            time_lapse = f"{end_time - start_time} seconds"

        except Exception as e:
            return JSONResponse(status_code=400, content={"error": str(e)})

        # NON GPT recommendation:
        keywords = extract_keywords(text=user_prompt, country_name=validated_country)
        try:
            scraped_imf = scrape_imf_search_results(keywords)
        except:
            scraped_imf = "Not found"

        try:
            scraped_wordlbank = scrape_worldbank_search_results(keywords)
        except:
            scraped_wordlbank = "Not found"

        result_concated = f"""{recommended_documents["result"]} \n\n<label style="color: #37377d">Recommendations based on online web sources:</label> \nSource: <a href="https://www.imf.org/" target="_blank">https://imf.org/</a>\n{scraped_imf} \n\nSource: <a href="https://www.worldbank.org/" target="_blank">https://www.worldbank.org/</a>\n{scraped_wordlbank}"""

        return {"recommended_documents": result_concated, "prompt":recommendation_template, "use_default_prompt":use_default_prompt, "country_name":validated_country, "is_super_admin":False}


@app.post("/generate_concept_note", status_code=200)
async def generate_concept_note(
        user_prompt: str,
        country_name: str,
        backend_context: Optional[str] = None,
        files: List[UploadFile] = File(None),
        custom_backend_prompt: Optional[str] = None,
        is_super_admin: Optional[bool] = False,
):
    """
    *****

    **********:
    - *****
    - *****
    - *****
    - *****
    - *****
         ** ***** **

            *****
            *****
            *****
            *****

            *****

            *****
            *****
            --------
            *****
            1.*****:
            *****

            2.*****:
            *****

            3.*****:
            *****

            4.*****:
            *****


    Returns:
    - dict: A dictionary containing the generated concept note and related information.
    """
    try:
        clean_context = ''
        if files:
            # read file as context
            context = read_files(files=files)
            # SUMMARIZE CONTEXT
            summarized_context = nltk_summarize(text=context,
                                                max_words=MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_OUTLINE)
            # clean up context
            clean_context += cleanup_text(summarized_context)
        if backend_context:
            # read streamed context from backend
            context = backend_context
            # SUMMARIZE CONTEXT
            summarized_context = nltk_summarize(text=context,
                                                max_words=MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_OUTLINE)
            # clean up context
            clean_context += cleanup_text(summarized_context)

        # masking context
        masker = Masker(clean_context)
        masked_context = masker.masked_data
        context = masked_context
        # count context token
        token_counts = count_tokens(masked_context)
        if token_counts > 100000:
            return JSONResponse(status_code=400,
                                content={"error": "Too many tokens, please upload fewer files."})
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})

    country_name = country_name
    user_prompt = user_prompt
    question = 'Concept Notes'
    # Concept VARIABELS

    # Admin has custom rules
    if is_super_admin == True:
        concept_chain = RetrievalQA.from_chain_type(
            llm=ChatOpenAI(
                model_name=MODEL_NAME,
                openai_api_key=OPENAI_API_KEY,
                max_tokens=MAX_TOKENS_COMPLETION_FOR_GENERATE_CONCEPT_NOTE,
                timeout=90,
            ),
            retriever=vectordb.as_retriever(search_kwargs={'k': 7}),
            return_source_documents=True,
        )

        use_default_prompt = True
        default_writing_rule = f""" *****
        
                                    ***** : {country_name} *****
                                    *****
                                    *****
                                    *****
                                    *****
                                    {context}
                                    
                                    
                                    *****
                                    ----
                                    *****
                                    *****:
                                    1. *****
                                    1.1 *****
                                            *****
                                    1.2 *****
                                            *****
                                    1.3 *****
                                            *****
                                    *****
                                    *****
                                    *****"""

        use_default_prompt = False
        concept_template = user_prompt + "\n-----------------\n" + default_writing_rule

        try:
            # Generate the policy outline using the context and question
            start_time = time.time()
            total_tokens = 0
            total_prompt_tokens = 0
            total_completion_tokens = 0
            total_cost = 0

            with get_openai_callback() as cb:
                concept_notes = await concept_chain.ainvoke(concept_template)
                concept_notes = concept_notes['result']
                concept_notes = masker.unmask_data(concept_notes)
                total_tokens = cb.total_tokens
                total_prompt_tokens = cb.prompt_tokens
                total_completion_tokens = cb.completion_tokens
                total_cost = cb.total_cost

            end_time = time.time()
            time_lapse = f"{end_time - start_time} seconds"
        except Exception as e:
            return JSONResponse(status_code=400, content={"error": str(e)})
        # Get title
        try:
            title = get_title(user_prompt=user_prompt, question=question, country_name=country_name)
        except Exception as e:
            return JSONResponse(status_code=400, content={"error": str(e)})

        # Parse result
        try:
            parsed_result = langchain_parsing_for_notes(concept_notes)
        except Exception as e:
            return JSONResponse(status_code=400, content={"error": str(e)})

        # return {"context": context, "country_name": country_name}
        return {
            "is_super_admin":is_super_admin,
            "title": title,
            "result": concept_notes,
            "parsed_result": parsed_result,
            "country_name": country_name,
            "use_default_prompt": use_default_prompt,
            "prompt": concept_template,
            "tokens": {
                "total_tokens": total_tokens,
                "prompt_tokens": total_prompt_tokens,
                "completion_prompt": total_completion_tokens,
                "cost": total_cost
            },
            "time_lapse": time_lapse
        }

    else:
        concept_chain = RetrievalQA.from_chain_type(
            llm=ChatOpenAI(
                model_name=MODEL_NAME,
                openai_api_key=OPENAI_API_KEY,
                max_tokens=MAX_TOKENS_COMPLETION_FOR_GENERATE_CONCEPT_NOTE,
                timeout=90,
            ),
            retriever=vectordb.as_retriever(search_kwargs={'k': 7}),
            return_source_documents=True,
        )

        use_default_prompt = True
        default_writing_rule =  f"""*****
                                    *****
                                    *****:
                                    1. *****
                                    1.1 *****
                                            *****
                                    1.2 *****
                                            *****
                                    1.3 *****
                                            *****
                                    *****
                                    *****
                                    *****"""

        if custom_backend_prompt:
            prompt_string = validate_prompt_input(custom_backend_prompt, PromptImplementation.CONCEPT_NOTES)
            prompt_string = prompt_string.replace('{user_prompt}', user_prompt)
            prompt_string = prompt_string.replace('{context}', context)
            prompt_string = prompt_string.replace('{country_name}', country_name)
            concept_template = prompt_string+"\n-----------------\n"+default_writing_rule
            # outline_template = outline_template.replace('{user_prompt}', user_prompt).replace('{country_name}', country_name).replace('{context}', context)
            use_default_prompt = False
        else:
            # country_name = detected_country_name
            concept_template = f"""
            *****, {user_prompt}. 
            *****, 
            ***** : {country_name} *****
            *****
            *****
    
            *****:
    
            *****
            *****
            *****
            *****
            *****: 
            {context}
            --------
            *****:
            1.*****:
            *****
    
            2.*****:
            *****
    
            3.*****:
            *****
    
            4.*****:
            *****
    
            --------------------------------------
            *****
            *****
            *****
            *****
            *****
            *****:
            1. *****
            1.1 *****
                    *****
            1.2 *****
                    *****
            1.3 *****
                    *****
            *****
            *****
            *****
            """
        try:
            # Generate the policy outline using the context and question
            start_time = time.time()
            total_tokens = 0
            total_prompt_tokens = 0
            total_completion_tokens = 0
            total_cost = 0

            with get_openai_callback() as cb:
                concept_notes = await concept_chain.ainvoke(concept_template)
                concept_notes = concept_notes['result']
                concept_notes = masker.unmask_data(concept_notes)
                total_tokens = cb.total_tokens
                total_prompt_tokens = cb.prompt_tokens
                total_completion_tokens = cb.completion_tokens
                total_cost = cb.total_cost

            end_time = time.time()
            time_lapse = f"{end_time - start_time} seconds"
        except Exception as e:
            return JSONResponse(status_code=400, content={"error": str(e)})
        # Get title
        try:
            title = get_title(user_prompt=user_prompt, question=question, country_name=country_name)
        except Exception as e:
            return JSONResponse(status_code=400, content={"error": str(e)})

        # Parse result
        try:
            parsed_result = langchain_parsing_for_notes(concept_notes)
        except Exception as e:
            return JSONResponse(status_code=400, content={"error": str(e)})

        # return {"context": context, "country_name": country_name}
        return {
            "is_super_admin": False,
            "title": title,
            "result": concept_notes,
            "parsed_result": parsed_result,
            "country_name": country_name,
            "use_default_prompt": use_default_prompt,
            "prompt": concept_template,
            "tokens": {
                "total_tokens": total_tokens,
                "prompt_tokens": total_prompt_tokens,
                "completion_prompt": total_completion_tokens,
                "cost": total_cost
            },
            "time_lapse": time_lapse
        }


@app.post("/generate_table_of_content", status_code=200)
async def generate_table_of_content(
        user_prompt: str,
        country_name: str,
        backend_context: Optional[str] = None,
        files: List[UploadFile] = File(None),
        custom_backend_prompt: Optional[str] = None,
):
    """
    *****

    **********:
    - *****
    - *****
    - *****
    - *****
    - *****
         ** Default prompt by system: **


            *****
            *****
            *****
            *****
            *****
            *****
            *****:
            --------
            *****:

            1. *****
            - *****
            - *****
            - *****

            2. *****
            - *****
            - *****
            - *****

            3. *****
            - *****
            - *****
            - *****

            4. *****
            - *****
            - *****
            - *****

            5. *****
            - *****
            - *****
            - *****

            6. *****
            - *****
            - *****
            - *****

            7. *****
            - *****
            - *****
            - *****

            8. *****
            - *****
            - *****
            - *****

            *****:
            - *****
            - *****
            - *****
            - *****
              *****:
               {context}
            - *****
            - *****
            - *****
            --------------------------------------
            *****
            *****
            *****
            *****
            *****:
            *****
            1. *****
                1.1 *****
            2. *****
                2.1 *****
                2.2 *****

            *****
            1. *****
                1.1 *****
                2.1 *****
            2. *****
                2.2 *****

            *****
            *****
            *****


    Returns:
    dict: A dictionary containing the generated table of contents and related information.

    Raises:
    Exception: If an error occurs during table of contents generation.
    """
    try:
        clean_context = ''
        if files:
            # read file as context
            context = read_files(files=files)
            # SUMMARIZE CONTEXT
            summarized_context = nltk_summarize(text=context,
                                                max_words=MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_OUTLINE)
            # clean up context
            clean_context += cleanup_text(summarized_context)
        if backend_context:
            # read streamed context from backend
            context = backend_context
            # SUMMARIZE CONTEXT
            summarized_context = nltk_summarize(text=context,
                                                max_words=MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_OUTLINE)
            # clean up context
            clean_context += cleanup_text(summarized_context)

        # masking context
        masker = Masker(clean_context)
        masked_context = masker.masked_data
        context = masked_context

        # count context token
        token_counts = count_tokens(masked_context)
        if token_counts > 100000:
            return JSONResponse(status_code=400,
                                content={"error": "Too many tokens, please upload fewer files."})
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})

    country_name = country_name
    user_prompt = user_prompt
    question = 'Policy Outline'

    # @TODO: limit user prompt token
    # @TODO: add if else from question (outline vs concept notes), or TBD to process at the same/different EP?
    # TOC VARIABELS
    use_default_prompt = True
    default_writing_rule = f""" *****
                                *****
                                *****
                                *****
                    
                                *****
                                1. *****
                                    *****
                                2. *****
                                    *****
                                    *****
                            
                                *****
                                1. *****
                                    *****
                                    *****
                                2. *****
                                    *****
                            
                                *****
                                *****
                                *****"""
    if custom_backend_prompt:
        prompt_string = validate_prompt_input(custom_backend_prompt, PromptImplementation.TABLE_OF_CONTENT)
        prompt_string = prompt_string.replace('{user_prompt}', user_prompt)
        prompt_string = prompt_string.replace('{country_name}', country_name)
        prompt_string = prompt_string.replace('{context}', context)
        toc_template = prompt_string + "\n-----------------\n" + default_writing_rule
        use_default_prompt = False
    else:
        toc_template = f"""
        *****, {user_prompt}. 
        *****, 
        ***** : {country_name}. 
        *****
        *****
        *****
        *****
        --------
        ***** :
    
        1. *****
        - *****
        - *****
        - *****
    
        2. *****
        - *****
        - *****
        - *****
    
        3. *****
        - *****
        - *****
        - *****
    
        4. *****
        - *****
        - *****
        - *****
    
        5. *****
        - *****
        - *****
        - *****
    
        6. *****
        - *****
        - *****
        - *****
    
        7. *****
        - *****
        - *****
        - *****
    
        8. *****
        - *****
        - *****
        - *****
    
        *****:
        - *****
        - *****
        - *****
        - *****
          *****
           {context}
        - *****
        - *****
        - *****
        --------------------------------------
        *****
        *****
        *****
        *****
        *****
        *****
        1. *****
            *****
        2. *****
            *****
            *****
    
        *****
        1. *****
            *****
            *****
        2. *****
            *****
    
        *****
        *****
        *****
        """

    outline_chain = RetrievalQA.from_chain_type(
        llm=ChatOpenAI(
            model_name=MODEL_NAME,
            openai_api_key=OPENAI_API_KEY,
            max_tokens=MAX_TOKENS_COMPLETION_FOR_GENERATE_OUTLINE,
            timeout=90,
        ),
        retriever=vectordb.as_retriever(search_kwargs={'k': 7}),
        return_source_documents=True,
    )

    try:
        # Generate the policy outline using the context and question
        start_time = time.time()
        total_tokens = 0
        total_prompt_tokens = 0
        total_completion_tokens = 0
        total_cost = 0

        with get_openai_callback() as cb:
            policy_outline = await outline_chain.ainvoke(toc_template)
            policy_outline = masker.unmask_data(policy_outline['result'])
            total_tokens = cb.total_tokens
            total_prompt_tokens = cb.prompt_tokens
            total_completion_tokens = cb.completion_tokens
            total_cost = cb.total_cost

        end_time = time.time()
        time_lapse = f"{end_time - start_time} seconds"
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})

    rouge_score = ''
    bert_score = ''
    try:
        start_time = time.time()
        scores = scorer.score(clean_context, policy_outline)
        for key in scores:
            rouge_score = f'{key}: {scores[key]}'
        end_time = time.time()
        scoring_time_lapse = f"{end_time - start_time} seconds"
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})
    # Get title
    try:
        title = get_title(user_prompt=user_prompt, question=question, country_name=country_name)
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})

    import json
    parsed_policy_outline = toc_parsing(table_of_content=policy_outline)

    return {
        "title": title,
        "result": policy_outline,
        "parsed_result": parsed_policy_outline,
        "country_name": country_name,
        "use_default_prompt": use_default_prompt,
        "prompt": toc_template,
        "tokens": {
            "total_tokens": total_tokens,
            "prompt_tokens": total_prompt_tokens,
            "completion_prompt": total_completion_tokens,
            "cost": total_cost
        },
        "rougeL_score": rouge_score,
        "bert_score": bert_score,
        "scoring_time_lapse": scoring_time_lapse,
        "chain_time_lapse": time_lapse
    }


@app.post("/regenerate_table_of_content", status_code=200)
async def regenerate_table_of_content(
        country_name: str,
        table_of_content: str,
        user_prompt: str,
        custom_backend_prompt: Optional[str] = None,
):
    """
    *****

    Parameters:
    - *****
    - *****
    - *****
    - *****
        ** Default prompt by system: **


        ***** :
        {table_of_content}
        -------------------
        *****
        {user_prompt}

        *****
        *****
        ***** {country_name}.
        *****


    Returns:
    dict: A dictionary containing the regenerated table of contents and related information.

    Raises:
    Exception: If an error occurs during table of contents regeneration.
    """

    # @TODO: limit user prompt token
    country_name = country_name
    table_of_content = table_of_content
    user_prompt = user_prompt

    context = 'Paraphrase / Regenerate'

    use_default_prompt = True
    default_writing_rule = f"""*****
                               *****
                               *****"""
    if custom_backend_prompt:
        prompt_string = validate_prompt_input(custom_backend_prompt, PromptImplementation.REGENERATE_TABLE_OF_CONTENT)
        prompt_string = prompt_string.replace('{table_of_content}', table_of_content)
        prompt_string = prompt_string.replace('{user_prompt}', user_prompt)
        prompt_string = prompt_string.replace('{country_name}', country_name)
        backend_prompt_for_regenerate = prompt_string + "\n-----------------\n" + default_writing_rule
        use_default_prompt = False
    else:
        backend_prompt_for_regenerate = f"""
        *****:
        {table_of_content}
        -------------------
        *****: 
        {user_prompt}
    
    
        *****
        *****
        ***** {country_name}. 
        *****
        """

    regenerate_outline_chain = RetrievalQA.from_chain_type(
        llm=ChatOpenAI(
            model_name=MODEL_NAME,
            openai_api_key=OPENAI_API_KEY,
            max_tokens=MAX_TOKENS_COMPLETION_FOR_GENERATE_OUTLINE,
            timeout=90,
            temperature=0,
        ),
        retriever=vectordb.as_retriever(search_kwargs={'k': 7}),
        return_source_documents=True
    )

    try:
        start_time = time.time()
        total_tokens = 0
        total_prompt_tokens = 0
        total_completion_tokens = 0
        total_cost = 0
        masker_prompt = Masker(backend_prompt_for_regenerate)

        with get_openai_callback() as cb:
            regenerated_result = await regenerate_outline_chain.ainvoke(masker_prompt.masked_data)
            regenerated_result = masker_prompt.unmask_data(regenerated_result['result'])
            total_tokens = cb.total_tokens
            total_prompt_tokens = cb.prompt_tokens
            total_completion_tokens = cb.completion_tokens
            total_cost = cb.total_cost

        parsed_result = toc_regenerated_parsing(regenerated_result)
        end_time = time.time()

        time_lapse = f"{end_time - start_time} seconds"
    except Exception as e:
        return JSONResponse(status_code=400,
                            content={"error": str(e)})

    return {
        "regenerated_result": regenerated_result,
        "parsed_result": parsed_result,
        "use_default_prompt": use_default_prompt,
        "country_name": country_name,
        "prompt": backend_prompt_for_regenerate,
        "tokens": {
            "total_tokens": total_tokens,
            "prompt_tokens": total_prompt_tokens,
            "completion_prompt": total_completion_tokens,
            "cost": total_cost
        },
        "time_lapse": time_lapse,
    }

@app.post("/generate_policy_outline", status_code=200)
async def generate_policy_outline(
        table_of_content: str,
        user_prompt: str,
        country_name: str,
        backend_context: Optional[str] = None,
        files: List[UploadFile] = File(None),
        custom_backend_prompt: Optional[str] = None,
):
    """
    *****

    *****:
    - *****
    - *****
    - *****
    - *****
    - *****
    - *****
        ** Default prompt by system: **


        *****
        *****
        *****
        *****
        *****
        *****
        ***** :
        --------
        ***** :

        1. *****
        - *****
        - *****
        - *****

        2. *****
        - *****
        - *****
        - *****

        3. *****
        - *****
        - *****
        - *****

        4. *****
        - *****
        - *****
        - *****

        5. *****
        - *****
        - *****
        - *****

        6. *****
        - *****
        - *****
        - *****

        7. *****
        - *****
        - *****
        - *****

        8. *****
        - *****
        - *****
        - *****

        *****:
        - *****
           {table_of_content}
        - *****
        - *****
        - *****
        - *****
          *****
           {context}
        - *****
        - *****
        - *****
        --------------------------------------
        *****
        *****
        *****
        *****
        *****
        1. *****
        1.1 *****
                *****
        1.2 *****
                *****
        *****
        *****


    Returns:
    - dict: A dictionary containing the generated policy outline, parsed result, and related information.
    """
    try:
        clean_context = ''
        if files:
            # read file as context
            context = read_files(files=files)
            # SUMMARIZE CONTEXT
            summarized_context = nltk_summarize(text=context,
                                                max_words=MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_OUTLINE)
            # clean up context
            clean_context += cleanup_text(summarized_context)
        if backend_context:
            # read streamed context from backend
            context = backend_context
            # SUMMARIZE CONTEXT
            summarized_context = nltk_summarize(text=context,
                                                max_words=MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_OUTLINE)
            # clean up context
            clean_context += cleanup_text(summarized_context)

        # masking context
        masker = Masker(clean_context)
        masked_context = masker.masked_data
        context = masked_context
        # count context token
        token_counts = count_tokens(masked_context)
        if token_counts > 100000:
            return JSONResponse(status_code=400,
                                content={"error": "Too many tokens, please upload fewer files."})
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})

    country_name = country_name
    user_prompt = user_prompt
    question = 'Policy Outline'
    # @TODO: limit user prompt token
    # @TODO: add if else from question (outline vs concept notes), or TBD to process at the same/different EP?


    # OUTLINE VARIABELS
    use_default_prompt = True
    default_writing_rule =  f"""*****
                                *****
                                *****
                                *****
                                *****
                                *****
                                1. *****
                                1.1 *****
                                        *****
                                1.2 *****
                                        *****
                                *****
                                *****"""
    if custom_backend_prompt:
        prompt_string = validate_prompt_input(custom_backend_prompt, PromptImplementation.POLICY_OUTLINE)
        prompt_string = prompt_string.replace('{user_prompt}', user_prompt)
        prompt_string = prompt_string.replace('{country_name}', country_name)
        prompt_string = prompt_string.replace('{context}', context)
        prompt_string = prompt_string.replace('{table_of_content}', table_of_content)
        outline_template = prompt_string+"\n-----------------\n"+default_writing_rule
        # outline_template = outline_template.replace('{user_prompt}', user_prompt).replace('{country_name}', country_name).replace('{context}', context)
        use_default_prompt = False
    else:
        outline_template = f"""
        *****, {user_prompt}. 
        *****, 
        ***** : {country_name}. 
        *****
        *****
        *****
        *****:
        --------
        *****:

        1. *****
        - *****
        - *****
        - *****

        2. *****
        - *****
        - *****
        - *****

        3. *****
        - *****
        - *****
        - *****

        4. *****
        - *****
        - *****
        - *****

        5. *****
        - *****
        - *****
        - *****

        6. *****
        - *****
        - *****
        - *****

        7. *****
        - *****
        - *****
        - *****

        8. *****
        - *****
        - *****
        - *****

        *****:
        - *****
           {table_of_content}
        - *****
        - *****
        - *****
        - *****
          *****
           {context}
        - *****
        - *****
        - *****
        --------------------------------------
        *****
        *****
        *****
        *****
        *****
        1. *****
        1.1 *****
                *****
        1.2 *****
                *****
        *****
        *****
        """

    outline_chain = RetrievalQA.from_chain_type(
        llm=ChatOpenAI(
            model_name=MODEL_NAME,
            openai_api_key=OPENAI_API_KEY,
            max_tokens=MAX_TOKENS_COMPLETION_FOR_GENERATE_OUTLINE,
            timeout=90,
        ),
        retriever=vectordb.as_retriever(search_kwargs={'k': 7}),
        return_source_documents=True,
    )

    try:
        # Generate the policy outline using the context and question
        start_time = time.time()
        total_tokens = 0
        total_prompt_tokens = 0
        total_completion_tokens = 0
        total_cost = 0

        with get_openai_callback() as cb:
            policy_outline = await outline_chain.ainvoke(outline_template)
            policy_outline = masker.unmask_data(policy_outline['result'])
            total_tokens = cb.total_tokens
            total_prompt_tokens = cb.prompt_tokens
            total_completion_tokens = cb.completion_tokens
            total_cost = cb.total_cost

        end_time = time.time()
        time_lapse = f"{end_time - start_time} seconds"
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})
    
    rouge_score = ''
    bert_score = ''
    try:
        start_time = time.time()
        scores = scorer.score(clean_context, policy_outline)
        for key in scores:
            rouge_score = f'{key}: {scores[key]}'

        end_time = time.time()
        scoring_time_lapse = f"{end_time - start_time} seconds"
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})
    # Get title
    try:
        title = get_title(user_prompt=user_prompt, question=question, country_name=country_name)
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})

    # Parse result
    try:
        parsed_result = langchain_parsing(policy_outline)
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})

    # return {"context": context, "country_name": country_name}
    return {
        "title": title,
        "result": policy_outline,
        "parsed_result": parsed_result,
        "country_name": country_name,
        "use_default_prompt": use_default_prompt,
        "prompt": outline_template,
        "tokens": {
            "total_tokens": total_tokens,
            "prompt_tokens": total_prompt_tokens,
            "completion_prompt": total_completion_tokens,
            "cost": total_cost
        },
        "rougeL_score": rouge_score,
        "bert_score": bert_score,
        "scoring_time_lapse": scoring_time_lapse,
        "chain_time_lapse": time_lapse
    }


async def process_subsection(section_title, subsections, updated_outline, context, country_name, backend_prompt,
                             use_default_prompt):
    # DETAILS VARIABELS
    details_chain = RetrievalQA.from_chain_type(
        llm=ChatOpenAI(
            model_name=MODEL_NAME,
            openai_api_key=OPENAI_API_KEY,
            max_tokens=MAX_TOKENS_COMPLETION_FOR_GENERATE_SUBSECTION,
            timeout=90,
        ),
        retriever=vectordb.as_retriever(search_kwargs={'k': 10}),
        return_source_documents=True,
    )
    updated_subsections = {}
    for subsection_title, content in subsections.items():
        # custom_backend_prompt
        if use_default_prompt == False:
            prompt_string = validate_prompt_input(backend_prompt, PromptImplementation.SUBSECTIONS_DETAILING)
            prompt_string = prompt_string.replace('{content}', content)
            prompt_string = prompt_string.replace('{subsection_title}', subsection_title)
            prompt_string = prompt_string.replace('{section_title}', section_title)
            prompt_string = prompt_string.replace('{context}', context)
            prompt_string = prompt_string.replace('{country_name}', country_name)
            details_template = prompt_string
        else:
            details_template = f"""
            *****:
            '{content}' 

            ***** '{subsection_title}', ***** '{section_title}'.
            *****
            *****

            *****
            {context}

            ------------
            *****
            *****

            ***** {country_name}. *****
            *****
            *****
            *****
            *****
            *****
            """
        invoke_result = await details_chain.ainvoke(details_template)
        updated_content = invoke_result['result']
        updated_subsections[subsection_title] = updated_content

    updated_outline[section_title] = updated_subsections


# UNHIDE IF CLIENT REQUEST
# @app.post("/generate_policy_subsections", status_code=200)
# async def generate_policy_subsections(
#         country_name: str,
#         policy_outline: str,
#         policy_title: str,
#         backend_context: Optional[str] = None,
#         files: list[UploadFile] = File(None),
#         custom_backend_prompt: Optional[str] = None,
# ):
#     try:
#         clean_context = ''
#         if files:
#             # read file as context
#             context = read_files(files=files)
#             # SUMMARIZE CONTEXT
#             summarized_context = nltk_summarize(text=context,
#                                                 max_words=MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_SUBSECTIONS)
#             # clean up context
#             clean_context += cleanup_text(summarized_context)
#         if backend_context:
#             # read file as context
#             context = backend_context
#             # SUMMARIZE CONTEXT
#             summarized_context = nltk_summarize(text=context,
#                                                 max_words=MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_SUBSECTIONS)
#             # clean up context
#             clean_context += cleanup_text(summarized_context)

#         # masking context
#         masker = Masker(clean_context)
#         masked_context = masker.masked_data
#         context = masked_context
#         # count context token
#         token_counts = count_tokens(masked_context)
#         if token_counts > 100000:
#             return JSONResponse(status_code=400,
#                                 content={"error": "Too many tokens, please upload fewer files."})
#     except Exception as e:
#         return JSONResponse(status_code=400, content={"error": str(e)})

#     country_name = country_name
#     policy_outline = policy_outline
#     policy_title = policy_title

#     use_default_prompt = True
#     if custom_backend_prompt:
#         use_default_prompt = False
#         default_details_template = custom_backend_prompt
#     else:
#         default_details_template = """
#         Please brief and/or expand and/or give answer this content:
#         {content}

#         For your information, that content is under the subsection {subsection_title}, in the section {section_title}.
#         The explanation should be concise, directly address the stated content, and relate to its section.
#         The provided explanation must be written in at least 3 paragraphs.

#         Also, please refer to any relevant country profile and/or strategies/problems mentioned to provide the specify answer here:
#         {context}

#         ------------
#         Ensure that the explanations are concise, valuable, directly address the mentioned content, and has specify answer based on the context.
#         If the content is already explained, you can expand the existing explanation.

#         The country mentioned is {country_name}. Do not mention other country names unless they are relevant to the explanation.
#         Do not mention problems/explanations outside the topic from the outline and context unless they are relevant.
#         Do not explain the title/position of the outline, or your role as a policy advisor, because at the moment you are detailing your document.
#         Do not mention again it is under a subsection or section.
#         Just to the point give the paragraphs of the result.
#         Just provide the answer/explanation.
#         """

#     try:
#         start_time = time.time()
#         total_tokens = 0
#         total_prompt_tokens = 0
#         total_completion_tokens = 0
#         total_cost = 0
#         import json
#         parsed_policy_outline = langchain_parsing(policy_outline=policy_outline)

#         updated_outline = {}
#         tasks = []
#         # Adding task to process each section asynchronously
#         for section_title, subsections in parsed_policy_outline.items():
#             task = process_subsection(section_title=section_title, subsections=subsections,
#                                       updated_outline=updated_outline, context=context, country_name=country_name,
#                                       backend_prompt=default_details_template, use_default_prompt=use_default_prompt)
#             tasks.append(task)

#         with get_openai_callback() as cb:
#             # run all task asynchronously
#             await asyncio.gather(*tasks)

#             total_tokens += cb.total_tokens
#             total_prompt_tokens += cb.prompt_tokens
#             total_completion_tokens += cb.completion_tokens
#             total_cost += cb.total_cost

#         end_time = time.time()
#         time_lapse = f"{end_time - start_time} seconds"
#     except Exception as e:
#         return JSONResponse(status_code=400, content={"error": str(e)})

#     try:
#         # sort result from async task
#         details_result = sort_dict_by_key(updated_outline)
#     except:
#         details_result = updated_outline

#     # return {"context": context, "country_name": country_name}
#     return {
#         "policy_outline": policy_outline,
#         "result": details_result,
#         "use_default_prompt": use_default_prompt,
#         "prompt": default_details_template,
#         "tokens": {
#             "total_tokens": total_tokens,
#             "prompt_tokens": total_prompt_tokens,
#             "completion_prompt": total_completion_tokens,
#             "cost": total_cost
#         },
#         "time_lapse": time_lapse
#     }


@app.post("/regenerate_subsection", status_code=200)
async def regenerate_subsection(
        country_name: str,
        subsection_title: str,
        subsection_explanation: str,
        user_prompt: str,
        custom_backend_prompt: Optional[str] = None,
):
    """
    *****

    Parameters:
    - *****
    - *****
    - *****
    - *****
    - *****
        ** Default prompt by system: **


        *****
        {subsection_title}
        *****
        {subsection_explanation}
        -------------------
        *****
        *****:
        {user_prompt_for_regenerate}

        ***** {country_name}.
        *****
        *****

        *****
        *****
        *****

        *****
        *****



    Returns:
    dict: A dictionary containing the regenerated subsection content and related information.

    Raises:
    Exception: If an error occurs during subsection regeneration.
    """

    # @TODO: limit user prompt token
    country_name = country_name
    subsection_title = subsection_title
    subsection_explanation = subsection_explanation
    user_prompt_for_regenerate = user_prompt

    validated_country = country_name
    if country_name.lower() == "other":
        try:
            detected_country_name = detect_country(sentence=user_prompt.lower())
            if detected_country_name is not None:
                validated_country = detected_country_name
            else:
                return JSONResponse(status_code=400, content={"error": "Country name in the user prompt not found"})
        except:
            return JSONResponse(status_code=400, content={"error": "Country name in the user prompt not found"})

    # DETAILS VARIABELS
    context = 'Paraphrase / Regenerate'

    use_default_prompt = True
    default_writing_rule = f"""*****
                            *****
                            *****
                            *****
                            
                            *****
                            *****"""
    if custom_backend_prompt:
        prompt_string = validate_prompt_input(custom_backend_prompt, PromptImplementation.REGENERATE_SUBSECTION)
        prompt_string = prompt_string.replace('{subsection_title}', subsection_title)
        prompt_string = prompt_string.replace('{subsection_explanation}', subsection_explanation)
        prompt_string = prompt_string.replace('{user_prompt_for_regenerate}', user_prompt_for_regenerate)
        prompt_string = prompt_string.replace('{country_name}', validated_country)
        backend_prompt_for_regenerate = prompt_string+"\n------------------\n"+default_writing_rule
        # outline_template = outline_template.replace('{user_prompt}', user_prompt).replace('{country_name}', country_name).replace('{context}', context)
        use_default_prompt = False
    else:
        backend_prompt_for_regenerate = f"""
        *****
        {subsection_title}
        *****
        {subsection_explanation}
        -------------------
        *****
        ***** :
        {user_prompt_for_regenerate}
        
        ***** {validated_country}. 
        *****
        *****
        
        *****
        *****
        *****

        *****
        *****
        """

    regenerate_chain = RetrievalQA.from_chain_type(
        llm=ChatOpenAI(
            model_name=MODEL_NAME,
            openai_api_key=OPENAI_API_KEY,
            max_tokens=MAX_TOKENS_COMPLETION_FOR_REGENERATE_SUBSECTION,
            timeout=90,
        ),
        retriever=vectordb.as_retriever(search_kwargs={'k': 7}),
        return_source_documents=True
    )

    try:
        start_time = time.time()
        total_tokens = 0
        total_prompt_tokens = 0
        total_completion_tokens = 0
        total_cost = 0
        masker_prompt = Masker(backend_prompt_for_regenerate)

        with get_openai_callback() as cb:
            regenerated_subsection = await regenerate_chain.ainvoke(masker_prompt.masked_data)
            regenerated_subsection = masker_prompt.unmask_data(regenerated_subsection['result'])

            total_tokens = cb.total_tokens
            total_prompt_tokens = cb.prompt_tokens
            total_completion_tokens = cb.completion_tokens
            total_cost = cb.total_cost

        end_time = time.time()
        time_lapse = f"{end_time - start_time} seconds"
    except Exception as e:
        return JSONResponse(status_code=400,
                            content={"error": str(e)})

    # return {"context": context, "country_name": country_name}
    return {
        "subsection_title": subsection_title,
        "regenerated_subsection_content": regenerated_subsection,
        "use_default_prompt": use_default_prompt,
        "country_name":validated_country,
        "prompt": backend_prompt_for_regenerate,
        "tokens": {
            "total_tokens": total_tokens,
            "prompt_tokens": total_prompt_tokens,
            "completion_prompt": total_completion_tokens,
            "cost": total_cost
        },
        "time_lapse": time_lapse
    }


# UNHIDE IF CLIENT REQUESTED
# @app.post("/regenerate_outline_or_notes", status_code=200)
# async def regenerate_outline_or_notes(
#         country_name: str,
#         doc_value: str,
#         doc_format: str = Query("Policy Outline", enum=["Policy Outline", "Concept Notes"]),
#         custom_backend_prompt: Optional[str] = None,
# ):
#     # @TODO: limit user prompt token
#     country_name = country_name
#     document_format = doc_format
#     document_value = doc_value
#     if doc_format == 'Policy Outline':
#         doc_max_tokens_suggested = MAX_TOKENS_COMPLETION_FOR_GENERATE_OUTLINE
#     else:
#         doc_max_tokens_suggested = MAX_TOKENS_COMPLETION_FOR_GENERATE_CONCEPT_NOTE

#     context = 'Paraphrase / Regenerate'

#     use_default_prompt = True
#     if custom_backend_prompt:
#         prompt_string = validate_prompt_input(custom_backend_prompt, PromptImplementation.REGENERATE_OUTLINES_NOTES)
#         prompt_string = prompt_string.replace('{document_format}', document_format)
#         prompt_string = prompt_string.replace('{document_value}', document_value)
#         prompt_string = prompt_string.replace('{country_name}', country_name)
#         backend_prompt_for_regenerate = prompt_string
#         # outline_template = outline_template.replace('{user_prompt}', user_prompt).replace('{country_name}', country_name).replace('{context}', context)
#         use_default_prompt = False
#     else:
#         backend_prompt_for_regenerate = f"""
#         Given the current {document_format}
#         {document_value}
#         -------------------
#         Please do a text regeneration or paraphase for the document above to have more subsections and more comprehensive explanations, 
#         But make sure only subsections number, title & it's value to be regenerated/paraphrased. Keep the sections number and name same. Keep the writing format like numbering, tabs same.

#         The country name is {country_name}. 
#         Do not mention other irrelevent country/topics unless there is relevancy to it.
#         Keep the result align with what discussed in previous subsection value.
#         """

#     regenerate_outline_chain = RetrievalQA.from_chain_type(
#         llm=ChatOpenAI(
#             model_name=MODEL_NAME,
#             openai_api_key=OPENAI_API_KEY,
#             max_tokens=doc_max_tokens_suggested,
#             timeout=90,
#             temperature=0,
#         ),
#         retriever=vectordb.as_retriever(search_kwargs={'k': 7}),
#         return_source_documents=True
#     )

#     try:
#         start_time = time.time()
#         total_tokens = 0
#         total_prompt_tokens = 0
#         total_completion_tokens = 0
#         total_cost = 0
#         masker_prompt = Masker(backend_prompt_for_regenerate)

#         with get_openai_callback() as cb:
#             regenerated_result = await regenerate_outline_chain.ainvoke(masker_prompt.masked_data)
#             regenerated_result = masker_prompt.unmask_data(regenerated_result['result'])
#             total_tokens = cb.total_tokens
#             total_prompt_tokens = cb.prompt_tokens
#             total_completion_tokens = cb.completion_tokens
#             total_cost = cb.total_cost

#         end_time = time.time()
#         time_lapse = f"{end_time - start_time} seconds"
#     except Exception as e:
#         return JSONResponse(status_code=400,
#                             content={"error": str(e)})

#     # return {"context": context, "country_name": country_name}
#     return {
#         "regenerated_result": regenerated_result,
#         "use_default_prompt": use_default_prompt,
#         "prompt": backend_prompt_for_regenerate,
#         "tokens": {
#             "total_tokens": total_tokens,
#             "prompt_tokens": total_prompt_tokens,
#             "completion_prompt": total_completion_tokens,
#             "cost": total_cost
#         },
#         "time_lapse": time_lapse,
#         "type": doc_format
#     }



@app.post("/convert_toc_to_string", status_code=200)
async def convert_toc_to_string(
        sections: dict
):
    result = ""
    sections = sections.get("sections", [])
    for section in sections:
        result += f"{section['section']}:\n"
        for subsection in section['subsections']:
            result += f" {subsection}\n"
        result += "\n"
    return JSONResponse(content=result)


@app.post("/generate_summarize", status_code=200)
async def generate_summarize(
        country_name: Optional[str] = None,
        backend_context: Optional[str] = None,
        files: List[UploadFile] = File(None),
):
    """
    *****

    *****:
    - *****
    - *****
    - *****

    Returns:
    - dict: A dictionary containing the summarized information.
    """
    try:
        clean_context = ''
        if files:
            # read file as context
            context = read_files(files=files)
            # SUMMARIZE CONTEXT
            summarized_context = nltk_summarize(text=context,
                                                max_words=MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_OUTLINE) #3000
            # clean up context
            clean_context += cleanup_text(summarized_context)
        if backend_context:
            # read streamed context from backend
            context = backend_context
            # SUMMARIZE CONTEXT
            summarized_context = nltk_summarize(text=context,
                                                max_words=MAX_WORDS_CONTEXT_SUMMARIZATION_FOR_GENERATE_OUTLINE) #3000
            # clean up context
            clean_context += cleanup_text(summarized_context)

        token_counted = count_tokens(clean_context)

    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})

    return {
        "result": clean_context,
        "token_counted_w_tiktoken": token_counted,
        "country_name": country_name,
    }
