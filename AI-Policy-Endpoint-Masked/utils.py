import re
import fitz
from docx import Document
from pptx import Presentation
import openpyxl
import io
import os
import shutil
from fastapi.responses import JSONResponse
from rake_nltk import Rake
from fastapi import FastAPI
from selenium import webdriver
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.common.by import By
from selenium.webdriver.support import expected_conditions as EC
from nltk.corpus import stopwords


def scrape_worldbank_search_results(query):
    query = query.replace(' ', '+')
    url = f'https://www.worldbank.org/en/search?q={query}'
    options = webdriver.ChromeOptions()
    options.add_argument("--headless")
    options.add_argument("--no-sandbox")
    options.add_argument("--disable-gpu")
    options.add_experimental_option('excludeSwitches', ['enable-logging'])
    options.add_argument('--ignore-certificate-errors')
    driver = webdriver.Chrome(options=options)

    try:
        driver.get(url)
        results = []
        # wait 15 seconds for all elements to be located
        elements = WebDriverWait(driver, 15).until(
            EC.visibility_of_all_elements_located(
                (By.CSS_SELECTOR, 'div.all__search_listingitem.all__result_group ul li a'))
        )

        # Extract titles, URLs, and descriptions from elements
        i = 0
        for element in elements:
            try:
                title_element = element.find_element(By.CSS_SELECTOR, 'h2.all__list_title')
                title = title_element.text.strip()
            except:
                title = ""

            try:
                url = element.get_attribute('href')
            except:
                url = ""

            i = i+1

            if title == "":
                results.append(f"{i}. <a href='{url}' target='_blank'>{url}</a>")
            else:
                results.append(f"{i}. <a href='{url}' target='_blank'>{title}</a>")

        if len(results)> 6:
            results = results[:5]
        return '\n'.join(results)

    finally:
        driver.quit()

def scrape_imf_search_results(query):
    url = f'https://www.imf.org/en/Search#q={query}&sort=relevancy&numberOfResults=50'
    options = webdriver.ChromeOptions()
    options.add_argument("--headless")
    options.add_argument("--no-sandbox")
    options.add_argument("--disable-gpu")
    options.add_experimental_option('excludeSwitches', ['enable-logging'])
    options.add_argument('--ignore-certificate-errors')
    driver = webdriver.Chrome(options=options)

    try:
        driver.get(url)
        elements = WebDriverWait(driver, 15).until(
            EC.visibility_of_all_elements_located((By.CSS_SELECTOR, '.imf-result-item .CoveoResultLink'))
        )
        results = []
        i = 0
        for element in elements:
            title = element.get_attribute('title')
            if (title is not None) and (title.strip() != ""):
                i = i+1
                url = element.get_attribute('href')
                results.append(f"{i}. <a href='{url}' target='_blank'>{title}</a>")
            else:
                i = i + 1
                url = element.get_attribute('href')
                results.append(f"{i}. <a href='{url}' target='_blank'>{url}</a>")


        results = results[:5]
        return '\n'.join(results)

    finally:
        driver.quit()


def extract_keywords(text, country_name):
    text = text.lower()
    r =  Rake().extract_keywords_from_text(text=text)
    country = f"in {country_name}"""

    keywords = None
    try:
        keywords = r.get_ranked_phrases()
        if country_name.lower() in keywords:
            keywords = keywords.remove(country_name.lower())
        if 'generate' in keywords:
            keywords = keywords.remove('generate')
        if 'please' in keywords:
            keywords = keywords.remove('please')
        keywords.append(country)
    except:
        if keywords is None:
            keywords = text.split()
    keywords_str = ' '.join(keywords)

    return keywords_str

def cleanup_text(text):
    # Remove html <a> tag
    text = re.sub(r"<a href[^>]*>([^<]+)</a>", " ", text)
    text = re.sub(r"<a rel[^>]*>([^<]+)</a>", " ", text)

    # Remove image-related tags
    text = re.sub(r"<img[^>]*>", " ", text)
    text = re.sub(r"<figure[^>]*>.*?</figure>", " ", text)
    # Remove image-related tags including .png extension
    text = re.sub(r"<img[^>]*>|<figure[^>]*>.*?</figure>|<[^>]*.png[^>]*>", " ", text)

    # Replace specific domain
    text = text.replace("WWW. QQGIAT .NET", " ")

    # Replace special characters
    text = text.replace("\t", " ").replace("\n", " ").replace("(\r", " ").replace("&nbsp;", " ").replace("amp;", " ")

    # Remove url link
    text = re.sub(r"http\S+", "", text)
    text = re.sub(r"www.\S+", "", text)

    # Keep letters and numbers only
    # text = re.sub(r"[^a-zA-Z0-9\s]", " ", text)
    text = re.sub(r"[^\w\s.,]", " ", text)

    # Keep single spaces
    text = re.sub(" +", " ", text)

    # Remove long sequences of periods, but keep single periods and other punctuation
    text = re.sub(r"(?<!\w)\.{3,}(?!\w)", " ", text)  # Replace 3 or more periods not surrounded by word characters
    text = re.sub(r"(?<!\w)\.{2,}(?!\w)", " ", text)  # Replace 2 or more periods not surrounded by word characters

    # Remove sequences of 3 or more periods
    text = re.sub(r"\.{3,}", " ", text)
    text = re.sub(r"\. {3,}", " ", text)

    # Keep single spaces
    text = re.sub(" +", " ", text)

    # Define the pattern to match unwanted characters (copied)
    pattern = re.compile(
        r"[\n\xe2\x96\xaa\xe2\x96\xaa\xe2\x80\x99\xe2\x80\x9c\xe2\x80\x9d\xe2\x80\x9d\xe2\x96\xba\xe2\x80\x99\xe2\x80\x99\xe2\x96\xba\xe2\x80\x9c\xe2\x80\x9d\xe2\x80\x9d\xe2\x80\x9c]")
    # Use the pattern to substitute the unwanted characters with an empty string (copied)
    text = re.sub(pattern, '', text)

    # remove long underscores
    text = re.sub(r'_+', '', text)

    return text


def read_files(files):
    concatenated_text = ""
    context = ""
    try:
        for file in files:
            context_file = file
            file_ext = context_file.filename.split(".")[-1]

            try:
                if file_ext == "txt":
                    context_bytes = context_file.file.read()
                    context = context_bytes.decode("utf-8")
                elif file_ext == "pdf":
                    # Read the file in memory
                    file_content = context_file.file.read()
                    # Extract text from the PDF content
                    text = ""
                    with fitz.open("pdf", file_content) as doc:
                        for page in doc:
                            text += page.get_text()
                    context = text
                elif file_ext == "docx":
                    context_bytes = context_file.file.read()
                    doc = Document(io.BytesIO(context_bytes))
                    context = ""
                    for paragraph in doc.paragraphs:
                        context += paragraph.text
                elif file_ext == "pptx":
                    context_bytes = context_file.file.read()
                    prs = Presentation(io.BytesIO(context_bytes))
                    context = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                context += shape.text + "\n"
                elif file_ext == "xlsx":
                    context_bytes = context_file.file.read()
                    wb = openpyxl.load_workbook(io.BytesIO(context_bytes))
                    context = ""
                    for sheet in wb.worksheets:
                        for row in sheet.iter_rows():
                            for cell in row:
                                context += str(cell.value) + "\n"
                elif file_ext == "doc":
                    # Get the current working directory
                    current_directory = os.getcwd()
                    # Define the path for the uploaded file
                    temp_folder = os.path.join(current_directory, 'temp_folder')
                    file_path = os.path.join(temp_folder, file.filename)

                    # Save the uploaded file
                    with open(file_path, "wb") as buffer:
                        shutil.copyfileobj(file.file, buffer)

                    output_file_path = os.path.join(temp_folder, os.path.splitext(file.filename)[0] + ".docx")

                    # Convert the file
                    import aspose.words as aw
                    doc = aw.Document(file_path)
                    doc.save(output_file_path)

                    # Read the context of the uploaded file
                    doc = Document(output_file_path)
                    context = ""
                    for paragraph in doc.paragraphs:
                        # remove watermark
                        if not paragraph.text.startswith("Evaluation Only. Created with Aspose.Words. Copyright"):
                            context += paragraph.text + "\n"
                        # Remove the uploaded file and the converted file
                    os.remove(file_path)
                    os.remove(output_file_path)
                elif file_ext == "ppt":
                    import aspose.slides as slides
                    # Get the current working directory
                    current_directory = os.getcwd()
                    # Define the path for the uploaded file
                    temp_folder = os.path.join(current_directory, 'temp_folder')
                    file_path = os.path.join(temp_folder, file.filename)

                    output_file_path = os.path.join(temp_folder, os.path.splitext(file.filename)[0] + ".pdf")

                    # Save the uploaded file
                    with open(file_path, "wb") as buffer:
                        shutil.copyfileobj(file.file, buffer)

                    presentation = slides.Presentation(file_path)

                    # Save the presentation as PDF
                    presentation.save(output_file_path, slides.export.SaveFormat.PDF)

                    # Read the file in memory
                    doc = fitz.open(output_file_path)
                    text = ""
                    for page in doc:
                        text += page.get_text()
                    context = text
                    os.remove(file_path)
                    os.remove(output_file_path)
                elif file_ext == 'xls':
                    import pandas as pd
                    import xlrd

                    # Read the uploaded file into memory
                    file_content = file.file.read()

                    # Open the Excel file using xlrd
                    workbook = xlrd.open_workbook(file_contents=file_content)
                    sh = workbook.sheet_by_index(0)

                    # Read data from each row and concatenate into context
                    context = ''
                    for rx in range(sh.nrows):
                        # Join each cell in the row into a single string
                        row_data = ' '.join(str(cell.value) for cell in sh.row(rx))
                        context += row_data + '\n'
                else:
                    return JSONResponse({"error": f"Unsupported file format: {file_ext}"})

                # SUMMARIZE CONTEXT
                # if count_tokens(context) > 3000:
                #     context = nltk_summarize(text=context, max_words=3000)
                concatenated_text += context + "\n"

            except Exception as e:
                return JSONResponse({"error": f"Failed to read the file: {e}"})

        return concatenated_text
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": f"Failed to read the file: {e}"})


def sort_dict_by_key(d):
    keys = list(d.keys())
    keys.sort()
    sorted_dict = {}
    for key in keys:
        sorted_dict[key] = d[key]
    return sorted_dict


def detect_country(sentence):
    country_list = [
      "Afghanistan", "Afghan", "United States", "United States of America", "USA",
      "Myanmar", "Burma", "Albania", "Algeria", "Andorra", "Angola", "Antigua and Barbuda",
      "Argentina", "Armenia", "Australia", "Austria", "Azerbaijan", "Bahamas", "Bahrain",
      "Bangladesh", "Barbados", "Belarus", "Belgium", "Belize", "Benin", "Bhutan",
      "Bolivia", "Bosnia and Herzegovina", "Botswana", "Brazil", "Brunei", "Bulgaria",
      "Burkina Faso", "Burundi", "Cabo Verde", "Cambodia", "Cameroon", "Canada",
      "Central African Republic", "Chad", "Chile", "China", "Colombia", "Comoros",
      "Congo", "Costa Rica", "Croatia", "Cuba", "Cyprus", "Czech Republic", "Denmark",
      "Djibouti", "Dominica", "Dominican Republic", "East Timor", "Ecuador", "Egypt",
      "El Salvador", "Equatorial Guinea", "Eritrea", "Estonia", "Eswatini", "Ethiopia",
      "Fiji", "Finland", "France", "Gabon", "Gambia", "Georgia", "Germany", "Ghana",
      "Greece", "Grenada", "Guatemala", "Guinea", "Guinea-Bissau", "Guyana", "Haiti",
      "Honduras", "Hungary", "Iceland", "India", "Indonesia", "Iran", "Iraq", "Ireland",
      "Israel", "Italy", "Jamaica", "Japan", "Jordan", "Kazakhstan", "Kenya", "Kiribati",
      "Korea, North", "Korea, South", "Kosovo", "Kuwait", "Kyrgyzstan", "Laos", "Latvia",
      "Lebanon", "Lesotho", "Liberia", "Libya", "Liechtenstein", "Lithuania", "Luxembourg",
      "Madagascar", "Malawi", "Malaysia", "Maldives", "Mali", "Malta", "Marshall Islands",
      "Mauritania", "Mauritius", "Mexico", "Micronesia", "Moldova", "Monaco", "Mongolia",
      "Montenegro", "Morocco", "Mozambique", "Namibia", "Nauru", "Nepal", "Netherlands",
      "New Zealand", "Nicaragua", "Niger", "Nigeria", "North Macedonia", "Norway", "Oman",
      "Pakistan", "Palau", "Palestine", "Panama", "Papua New Guinea", "Paraguay", "Peru",
      "Philippines", "Poland", "Portugal", "Qatar", "Romania", "Russia", "Rwanda",
      "Saint Kitts and Nevis", "Saint Lucia", "Saint Vincent and the Grenadines", "Samoa",
      "San Marino", "Sao Tome and Principe", "Saudi Arabia", "Senegal", "Serbia",
      "Seychelles", "Sierra Leone", "Singapore", "Slovakia", "Slovenia", "Solomon Islands",
      "Somalia", "South Africa", "South Sudan", "Spain", "Sri Lanka", "Sudan", "Suriname",
      "Sweden", "Switzerland", "Syria", "Taiwan", "Tajikistan", "Tanzania", "Thailand",
      "Togo", "Tonga", "Trinidad and Tobago", "Tunisia", "Turkey", "Turkmenistan", "Tuvalu",
      "Uganda", "Ukraine", "United Arab Emirates", "United Kingdom", "Uruguay", "Uzbekistan",
      "Vanuatu", "Vatican City", "Venezuela", "Vietnam", "Yemen", "Zambia", "Zimbabwe"
    ]

    sentence_lower = sentence.lower()

    # Iterate through each country name
    for country in country_list:
        # Lowercase the country name for case-insensitive comparison
        country_lower = country.lower()

        # Check if the country name is in the sentence
        if country_lower in sentence_lower:
            return country_lower

    # If no country name is found in the sentence
    return None
